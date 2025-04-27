# slide_generator_modified.py
import os
import io
import json
import time
import requests
import streamlit as st
from groq import Groq
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from urllib.parse import quote
import re
import base64 # Needed for local file handling

# Removed imports for block diagram generators as requested.
# Assuming PlantText and Manim are no longer used.
# PLANTTEXT_AVAILABLE = False # Explicitly mark as unavailable
# MANIM_AVAILABLE = False # Explicitly mark as unavailable


class SlideGenerator:
    """Core class for generating presentation slides using Groq LLM and Pollinations images."""

    def __init__(self):
        """Initialize the SlideGenerator with API clients and configurations."""
        self.groq_client = None
        # Removed diagram_generator initialization
        self.model = os.environ.get("GROQ_MODEL", "llama-3.1-70b-versatile") # Consider Llama 3.1

        # Removed configuration flags for diagram generators

        # Initialize Groq LLM client
        try:
            api_key = os.environ.get("GROQ_API_KEY")
            if api_key:
                self.groq_client = Groq(api_key=api_key)
                st.write("SlideGenerator: Initialized Groq client.") # Debug
            else:
                st.error("SlideGenerator: Missing GROQ_API_KEY environment variable. Groq features disabled.")
        except Exception as e:
            st.error(f"SlideGenerator: Failed to initialize Groq client: {e}")

        # Base URL for Pollinations image generation
        self.pollinations_base_url = os.environ.get("POLLINATIONS_API_URL", "https://image.pollinations.ai/prompt/")
        self.placeholder_image_url = "https://placehold.co/1024x576/eee/ccc?text=Image+Gen+Error"


    def parse_json_from_response(self, response_text):
        """
        Extract JSON from a string potentially containing JSON with START/END markers,
        handling potential control characters and providing detailed error logging.
        (Robust parsing logic - kept from original)
        """
        json_text = None
        parsed_json = {} # Default to empty dict

        try:
            # 1. Try extracting content between START OF JSON and END OF JSON markers
            json_match = re.search(r'START OF JSON\s*([\s\S]*?)\s*END OF JSON', response_text, re.DOTALL)
            if json_match:
                json_text = json_match.group(1).strip()
            else:
                # 2. If markers aren't found, assume the whole response might be JSON
                stripped_response = response_text.strip()
                if stripped_response.startswith('{') or stripped_response.startswith('['):
                    json_text = stripped_response
                else:
                     # 3. Try finding the first JSON object or array pattern if markers fail
                    object_match = re.search(r'({.*})', response_text, re.DOTALL)
                    if object_match:
                        json_text = object_match.group(1).strip()
                    else:
                        array_match = re.search(r'(\[[\s\S]*\])', response_text, re.DOTALL) # More robust array matching
                        if array_match:
                            json_text = array_match.group(1).strip()

            # 4. If we extracted potential JSON text, try parsing it
            if json_text:
                try:
                    parsed_json = json.loads(json_text)
                except json.JSONDecodeError as e1:
                    st.warning(f"Standard JSON parsing failed: {e1}. Trying with strict=False...")
                    try:
                        parsed_json = json.loads(json_text, strict=False)
                    except json.JSONDecodeError as e2:
                        error_pos = getattr(e2, 'pos', 0)
                        context_window = 30
                        start = max(0, error_pos - context_window)
                        end = min(len(json_text), error_pos + context_window)
                        error_context = json_text[start:end]
                        printable_context = ''.join(c if c.isprintable() or c in '\n\r\t' else f'\\x{ord(c):02x}' for c in error_context)

                        st.error(f"Failed to parse JSON even with strict=False: {e2}")
                        st.error(f"Error type: {type(e2).__name__}, Message: {e2.msg}")
                        st.error(f"Error at character {error_pos}: ...{printable_context}...")
                        st.error(f"Full problematic text snippet (around error):\n---\n{json_text[max(0, error_pos-150):min(len(json_text), error_pos+150)]}\n---")
                        parsed_json = {}

            else:
                st.error("Could not find any JSON structure in the response.")
                st.error(f"Full response received:\n{response_text}")
                parsed_json = {}

        except Exception as e:
            st.error(f"An unexpected error occurred during JSON parsing: {str(e)}")
            st.error(f"Problematic response text snippet:\n{response_text[:500]}...")
            parsed_json = {}

        if not isinstance(parsed_json, (dict, list)):
             st.warning(f"Parsing resulted in unexpected type {type(parsed_json)}, returning empty dict.")
             return {}

        return parsed_json


    def init_groq_client_if_needed(self):
        """Initialize the Groq client if it hasn't been initialized yet."""
        # This method is largely redundant now as initialization is in __init__
        if self.groq_client is None:
            api_key = os.environ.get("GROQ_API_KEY")
            if api_key:
                st.write("Initializing Groq client (on demand)...") # Debug
                try:
                    self.groq_client = Groq(api_key=api_key)
                     # Removed diagram generator re-initialization
                except Exception as e:
                     st.error(f"Failed to initialize Groq client on demand: {e}")
                     raise ValueError(f"Failed to initialize Groq client: {e}")
            else:
                st.error("Missing GROQ_API_KEY environment variable")
                raise ValueError("Missing GROQ_API_KEY environment variable")


    def generate_slide_structure(self, topic, delivery_medium, complexity_level, reference_text=None):
        """
        Generate the overall slide structure based on a topic or reference material.
        Only includes "image_slide" for visual representation where appropriate.
        """
        try:
            self.init_groq_client_if_needed() # Ensure client is ready

            # --- MODIFIED PROMPT ---
            # Removed block_diagram_slide type and related instructions.
            # Emphasized using image_slide for visual concepts within Pollinations' capability (general images).
            prompt = f"""
            Create a well-structured lecture slide deck outline on: "{topic}"
            The lecture is for {complexity_level.lower()} level students, delivered {delivery_medium}.

            Your response MUST follow these rules STRICTLY:
            1. Provide ONLY JSON output strictly between START OF JSON and END OF JSON markers. Do NOT include any text before START OF JSON or after END OF JSON.
            2. The JSON MUST be a valid list of JSON objects.
            3. Each object MUST represent a slide and contain 'slide_number' (int), 'slide_title' (string), and 'slide_type' (string).
            4. Use ONLY these exact slide types: "title_slide", "content_slide", "bullet_point_slide", "image_slide", "conclusion_slide".
            5. Use "image_slide" for general illustrative images (generated by Pollinations) that help visualize a concept. Consider topics that can be easily represented by a single, clear image.
            6. Ensure all strings in the JSON use double quotes ("") and any special characters within strings are correctly escaped (e.g., \\n, \\").
            7. IMPORTANT: If an "image_slide" is required to explain a topic, consider if a preceding "content_slide" with the same "slide_title" is needed to explain the concept verbally first. Pair them logically.

            START OF JSON
            [
                {{
                    "slide_number": 1,
                    "slide_title": "Introduction to {topic}",
                    "slide_type": "title_slide"
                }},
                {{
                    "slide_number": 2,
                    "slide_title": "[Specific Content-Driven Title]",
                    "slide_type": "content_slide"
                }},
                 {{
                    "slide_number": 3,
                    "slide_title": "[Concept Needing Visual Example]",
                    "slide_type": "content_slide"
                }},
                {{
                    "slide_number": 4,
                    "slide_title": "[Concept Needing Visual Example]",
                    "slide_type": "image_slide"
                }},
                ... more slides based on topic and complexity ...
                 {{
                    "slide_number": N,
                    "slide_title": "Conclusion",
                    "slide_type": "conclusion_slide"
                }}
            ]
            END OF JSON

            Decide the total number of slides based on {complexity_level.lower()} complexity and the depth needed for "{topic}".
            Make titles specific. Use "image_slide" for illustrations of concepts that Pollinations AI can generally depict (e.g., objects, scenes, abstract representations).
            Create paired content/visual slides where appropriate. Generate valid JSON according to the rules above.
            """
            # --- END OF MODIFIED PROMPT ---

            st.write("Calling Groq API for slide structure...") # Debug
            response = self.groq_client.chat.completions.create(
                messages=[{"role": "user", "content": prompt}],
                model=self.model,
                max_tokens=2500, # Increased slightly for potentially more complex structures
                temperature=0.7
            )
            result_text = response.choices[0].message.content
            st.write("Received structure response from Groq.") # Debug

            slide_structure = self.parse_json_from_response(result_text)

            if not isinstance(slide_structure, list) or not slide_structure:
                st.error("Failed to generate or parse slide structure. The LLM response might be invalid or empty.")
                if not isinstance(slide_structure, (dict, list)):
                     st.text_area("Raw Failed Structure Response", result_text, height=150)
                return []

            processed_structure = []
            # Removed block_diagram_slide from valid types
            valid_slide_types = {"title_slide", "content_slide", "bullet_point_slide", "image_slide", "conclusion_slide"}
            for i, slide in enumerate(slide_structure):
                 if isinstance(slide, dict) and "slide_title" in slide and "slide_type" in slide:
                    slide_type = slide.get("slide_type")
                    if slide_type not in valid_slide_types:
                        st.warning(f"Slide {i+1} ('{slide.get('slide_title')}') has invalid type '{slide_type}'. Replacing with 'content_slide'.")
                        slide["slide_type"] = "content_slide" # Default to content slide

                    slide["slide_number"] = i + 1 # Ensure sequential numbering
                    processed_structure.append(slide)
                 else:
                    st.warning(f"Skipping invalid slide structure item at index {i}: {slide}")

            st.session_state['debug_slide_structure'] = processed_structure
            st.write(f"DEBUG: Processed structure has {len(processed_structure)} slides.") # Debug

            return processed_structure

        except Exception as e:
            st.error(f"Error generating slide structure: {str(e)}")
            if 'result_text' in locals():
                st.text_area("Raw Response during Structure Exception", result_text, height=100)
            return []


    def generate_slide_content(self, slide_info, topic, complexity_level):
        """
        Generate detailed content for a specific slide, handling different types.
        Block diagram handling removed.
        """
        try:
            self.init_groq_client_if_needed() # Ensure client is ready

            slide_title = slide_info.get("slide_title", "Untitled Slide")
            slide_type = slide_info.get("slide_type", "content_slide")
            slide_number = slide_info.get("slide_number", 0)

            # --- MODIFIED PROMPT GENERATION LOGIC ---
            # Removed specific handling for block_diagram_slide

            json_rules = """
            Your response MUST follow these rules STRICTLY:
            1. Provide ONLY JSON output strictly between START OF JSON and END OF JSON markers. Do NOT include any text before START OF JSON or after END OF JSON.
            2. The JSON MUST be a valid JSON object.
            3. Ensure all strings in the JSON use double quotes ("") and escape special characters if needed (e.g., \\n, \\").
            """

            prompt = "" # Initialize prompt string

            if slide_type == "title_slide":
                prompt = f"""
                Create content for the title slide of a {complexity_level.lower()} lecture on "{topic}".
                Slide title: "{slide_title}"
                {json_rules}
                START OF JSON
                {{
                    "title": "{slide_title}",
                    "subtitle": "An Engaging Subtitle Relevant to {topic}",
                    "presenter": "Your Name / Affiliation",
                    "date": "{time.strftime('%B %d, %Y')}"
                }}
                END OF JSON
                """
            elif slide_type == "content_slide":
                # Check if this slide precedes an image slide with the same title
                is_explanation_for_image = False
                if 'debug_slide_structure' in st.session_state and slide_number < len(st.session_state['debug_slide_structure']):
                    # Correcting index: slide_number is 1-based, list index is 0-based. Next slide is at index `slide_number`.
                    if slide_number < len(st.session_state['debug_slide_structure']):
                        next_slide_info = st.session_state['debug_slide_structure'][slide_number]
                        if next_slide_info.get("slide_title") == slide_title and next_slide_info.get("slide_type") == "image_slide":
                            is_explanation_for_image = True


                explanation_guidance = "Explain the core concept clearly and concisely."
                if is_explanation_for_image:
                    explanation_guidance = f"This slide explains the concept that will be visualized in the *next* image slide. Focus on setting up the visual explanation. Keep text concise."

                prompt = f"""
                Create concise text content for slide {slide_number}, titled "{slide_title}" (type: {slide_type}), for a {complexity_level.lower()} lecture on "{topic}".
                {explanation_guidance} Aim for 2-3 short paragraphs or key points.
                {json_rules}
                START OF JSON
                {{
                    "main_content": [
                        "Paragraph 1: Introduce the core idea of '{slide_title}'.",
                        "Paragraph 2: Elaborate with key details or context relevant to {topic}.",
                        "Paragraph 3 (Optional): Add significance or a brief example."
                    ],
                    "needs_image": false
                }}
                END OF JSON
                """
            elif slide_type == "bullet_point_slide":
                prompt = f"""
                Create 3-6 concise bullet points for slide {slide_number}, titled "{slide_title}" (type: {slide_type}), for a {complexity_level.lower()} lecture on "{topic}".
                Each point should be short and impactful (under ~150 characters).
                {json_rules}
                START OF JSON
                {{
                    "main_content": [
                        "Key point 1 about '{slide_title}'.",
                        "Key point 2, distinct and clear.",
                        "Key point 3, adding another facet.",
                        "Key point 4 (if necessary).",
                        "Key point 5 (if necessary)."
                    ],
                    "needs_image": false
                }}
                END OF JSON
                """
            elif slide_type == "image_slide":
                 # Generate description for Pollinations AI
                 prompt = f"""
                 Create content for slide {slide_number}, titled "{slide_title}" (type: {slide_type}), for a {complexity_level.lower()} lecture on "{topic}".
                 This slide needs a general illustrative image from Pollinations.ai.
                 Generate a descriptive prompt for Pollinations.ai to create a relevant and professional image.
                 Also include a short caption for the slide. If any text is included in image strictly make sure the spelling is correct.
                 {json_rules}
                 START OF JSON
                 {{
                     "main_content": [
                         "Caption: Visual illustration related to {slide_title}"
                     ],
                     "needs_image": true,
                     "image_description": "A clear, professional illustration depicting [concept related to '{slide_title}'] for a {complexity_level.lower()} lecture on '{topic}'. Style: educational, clean background, high resolution."
                 }}
                 END OF JSON
                 """

            elif slide_type == "conclusion_slide":
                prompt = f"""
                Create content for the final conclusion slide {slide_number}, titled "{slide_title}" (type: {slide_type}), for a {complexity_level.lower()} lecture on "{topic}".
                Summarize 2-4 key takeaways concisely.
                {json_rules}
                START OF JSON
                {{
                    "main_content": [
                        "Key Takeaway 1: [Summarize main point from lecture]",
                        "Key Takeaway 2: [Summarize another crucial concept]",
                        "Key Takeaway 3: [Summarize application or implication]",
                        "Next Steps / Questions?"
                    ],
                    "needs_image": false
                }}
                END OF JSON
                """
            else: # Fallback for unknown type
                st.warning(f"Unknown slide type '{slide_type}' encountered for slide {slide_number}. Using generic content.")
                prompt = f"""
                Create generic content for slide {slide_number}, titled "{slide_title}" (type: {slide_type}), for a {complexity_level.lower()} lecture on "{topic}".
                {json_rules}
                START OF JSON
                {{
                    "main_content": [
                        "Placeholder content for '{slide_title}'.",
                        "This slide type ('{slide_type}') needs specific handling."
                    ],
                    "needs_image": False # Assuming unknown type doesn't need image by default
                }}
                END OF JSON
                """
            # --- END OF MODIFIED PROMPT GENERATION LOGIC ---


            # Call the Groq LLM API for content generation
            st.write(f"Calling Groq API for content of slide {slide_number} ('{slide_title}', Type: {slide_type})...") # Debug
            response = self.groq_client.chat.completions.create(
                messages=[{"role": "user", "content": prompt}],
                model=self.model,
                max_tokens=1000, # Adjust as needed
                temperature=0.7
            )
            result_text = response.choices[0].message.content
            st.write(f"Received content response for slide {slide_number}.") # Debug

            content = self.parse_json_from_response(result_text)

            if not isinstance(content, dict) or not content:
                st.warning(f"Content generation or parsing failed for slide '{slide_title}'. Using placeholder content.")
                # Add an error flag or message to the content dict
                content = {"main_content": ["Error: Could not generate content."], "needs_image": False, "error": "Failed to generate or parse content"}

            # Handle image generation for 'image_slide' using Pollinations
            if content.get("needs_image") and slide_type == "image_slide":
                image_description = content.get("image_description", f"Illustration for {slide_title}")
                st.write(f"Generating Pollinations image for slide {slide_number}...") # Debug
                image_url = self.generate_pollinations_image(image_description, topic)
                content["image_url"] = image_url

            # Debug output
            if 'debug_slide_contents' not in st.session_state: st.session_state['debug_slide_contents'] = []
            st.session_state['debug_slide_contents'].append({"slide_number": slide_number, **content})

            return content

        except Exception as e:
            st.error(f"Error generating content for slide '{slide_title}': {str(e)}")
            if 'result_text' in locals():
                st.text_area(f"Raw Response during Content Exception (Slide {slide_number})", result_text, height=100)
            return {"main_content": [f"Error generating content: {str(e)}"], "needs_image": False, "error": str(e)}


    def generate_pollinations_image(self, description, context):
        """
        Generate an image URL using the Pollinations API.
        """
        try:
            max_desc_len = 300
            if len(description) > max_desc_len:
                description = description[:max_desc_len] + "..."

            # More generic prompt for illustrations
            image_prompt = f"Educational illustration: {description}. Context: {context}. Style: clean, professional, clear, high resolution."

            encoded_prompt = quote(image_prompt)
            image_url = f"{self.pollinations_base_url}{encoded_prompt}"

            cache_buster = int(time.time())
            width = 1024
            height = 576 # 16:9 aspect ratio
            image_url += f"?width={width}&height={height}&seed={cache_buster}&nologo=true"

            st.write(f"DEBUG: Generated Pollinations Image URL: {image_url}") # Debug output
            return image_url

        except Exception as e:
            st.error(f"Error generating Pollinations image URL: {str(e)}")
            return self.placeholder_image_url # Return placeholder


    def create_slides_with_content(self, topic, delivery_medium, complexity_level, reference_text=None):
        """
        Generate all slides with their content in one coordinated process.
        """
        st.info("1. Generating slide structure...")
        # Clear previous debug info if any
        if 'debug_slide_structure' in st.session_state: del st.session_state['debug_slide_structure']
        if 'debug_slide_contents' in st.session_state: del st.session_state['debug_slide_contents']

        slide_structure = self.generate_slide_structure(
            topic, delivery_medium, complexity_level, reference_text
        )

        if not slide_structure:
            st.error("Failed to generate slide structure. Cannot proceed.")
            return [], []

        st.success(f"Generated structure with {len(slide_structure)} slides.")
        st.json(st.session_state.get('debug_slide_structure', []), expanded=False) # Show structure for debug

        st.info(f"2. Generating content for {len(slide_structure)} slides...")

        slide_contents = [None] * len(slide_structure)
        progress_bar = st.progress(0, text="Generating slide content...")
        total_slides = len(slide_structure)

        for i, slide_info in enumerate(slide_structure):
            slide_num = slide_info.get("slide_number", i + 1)
            slide_title = slide_info.get("slide_title", f"Slide {slide_num}")
            slide_type = slide_info.get("slide_type", "content_slide")

            progress_text = f"Generating: Slide {slide_num}/{total_slides} ('{slide_title}', Type: {slide_type})"
            progress_bar.progress((i + 1) / total_slides, text=progress_text)
            st.write(progress_text) # Show current step in main area too

            content = self.generate_slide_content(slide_info, topic, complexity_level)

            if 1 <= slide_num <= total_slides:
                 slide_contents[slide_num - 1] = content
            else:
                 st.error(f"Invalid slide number {slide_num} encountered. Content may be misplaced.")
                 if i < len(slide_contents): slide_contents[i] = content
                 else: slide_contents.append(content) # Append if index is totally wrong


        progress_bar.empty()
        final_contents = [c for c in slide_contents if c is not None]

        if len(final_contents) == total_slides:
            st.success(f"Successfully generated content for all {total_slides} slides.")
        else:
             st.warning(f"Generated content for {len(final_contents)} out of {total_slides} slides. Some may be missing due to errors.")

        st.write("Debug: Final Slide Contents:") # Show final contents for debug
        st.json(st.session_state.get('debug_slide_contents', []), expanded=False)

        # Removed Manim temporary file cleanup
        # if self.USE_MANIM_FOR_DIAGRAMS and isinstance(self.diagram_generator, ManimDiagramGenerator):
        #     self.diagram_generator.cleanup_temp_dir()


        return slide_structure, final_contents


    def create_powerpoint(self, slide_structure, slide_contents, topic):
        """
        Create a PowerPoint presentation from generated slide structure and content.
        Handles 'image_slide' from URL. Block diagram handling removed.
        """
        prs = Presentation()
        prs.slide_width = Inches(13.33) # 16:9 aspect ratio
        prs.slide_height = Inches(7.5)

        # Define layout indices (standard PowerPoint indices)
        # 0: Title Slide
        # 1: Title and Content
        # 5: Title Only
        # 6: Blank  <-- **NOTE: Standard index for Blank is often 6**
        # Verify these indices with your default PowerPoint template if issues persist.
        layout_indices = {
            'title_slide': 0,
            'content_slide': 1,
            'bullet_point_slide': 1,
            'image_slide': 6,        # Use Blank layout (index 6)
            'conclusion_slide': 1,
            'title_only': 5 # Index 5 is typically Title Only
            # Add other mappings if needed
        }
        # Default layout index if type not found
        default_layout_index = 1 # Title and Content

        st.info(f"3. Creating PowerPoint file with {len(slide_structure)} slides...")
        progress_bar = st.progress(0, text="Creating PowerPoint slides...")
        total_slides = len(slide_structure)

        if len(slide_structure) != len(slide_contents):
            st.error(f"Mismatch between slide structure ({len(slide_structure)}) and content ({len(slide_contents)}). Presentation may be incomplete.")
            num_slides_to_create = min(len(slide_structure), len(slide_contents))
        else:
            num_slides_to_create = total_slides

        for i in range(num_slides_to_create):
            progress_text = f"Creating slide {i+1}/{total_slides}"
            progress_bar.progress((i + 1) / total_slides, text=progress_text)

            slide_info = slide_structure[i]
            content = slide_contents[i]

            if not content or content.get("error"):
                st.warning(f"Skipping slide {i+1} ('{slide_info.get('slide_title', 'N/A')}') due to missing or error content.")
                continue

            slide_type = slide_info.get("slide_type", "content_slide")
            slide_title = slide_info.get("slide_title", f"Slide {i+1}")

            # Get the layout index for the current slide type
            layout_index = layout_indices.get(slide_type, default_layout_index)

            # Get the actual layout object from the presentation's layouts
            try:
                slide_layout = prs.slide_layouts[layout_index]
            except IndexError:
                 st.warning(f"Layout index {layout_index} not found in presentation template for slide type '{slide_type}'. Using default layout index {default_layout_index}.")
                 slide_layout = prs.slide_layouts[default_layout_index]

            # Add the slide using the chosen layout
            slide = prs.slides.add_slide(slide_layout)

            # --- SLIDE CREATION LOGIC ---
            # Special handling for image slides using blank layout (index 6)
            if slide_type == "image_slide":
                # Since we used the Blank layout (index 6), we add title manually
                # Check if a title placeholder exists on this specific blank layout (some might have one)
                title_placeholder = None
                for shape in slide.placeholders:
                     # Check for common title placeholder types or index 0
                    if shape.is_placeholder and (shape.placeholder_format.type == 1 or shape.placeholder_format.idx == 0):
                        title_placeholder = shape
                        break

                if title_placeholder:
                    title_shape = title_placeholder
                    title_shape.text = slide_title
                    self._format_text_shape(title_shape, font_size=36, bold=True) # Apply formatting
                     # Position title properly if placeholder isn't standard
                    if title_placeholder.placeholder_format.type != 1: # If it's not the standard title placeholder
                         title_shape.left = Inches(0.5)
                         title_shape.top = Inches(0.2)
                         title_shape.width = Inches(12.33)
                         title_shape.height = Inches(0.8)
                         self._format_text_shape(title_shape, font_size=36, bold=True, alignment=PP_ALIGN.CENTER) # Re-apply formatting with centering

                else:
                     # Add textbox for title if no title placeholder
                     title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.8))
                     title_shape.text = slide_title
                     self._format_text_shape(title_shape, font_size=36, bold=True, alignment=PP_ALIGN.CENTER)


                # Determine image source: only URL (Pollinations) is supported now
                image_source = None
                if "image_url" in content and content["image_url"]:
                     image_source = {"type": "url", "url": content["image_url"]}
                elif "image_content_base64" in content and content["image_content_base64"]:
                     # Fallback: if a base64 placeholder was generated due to generator error (less likely now)
                     image_source = {"type": "base64", "data": content["image_content_base64"]}


                if image_source:
                    try:
                        img_width_in = 8.0 # Target image width
                        img_aspect_ratio = 16.0 / 9.0 # Standard slide aspect ratio
                        img_height_in = img_width_in / img_aspect_ratio
                        img_left = (prs.slide_width.inches - img_width_in) / 2 # Center horizontally
                        img_top = Inches(1.2) # Position below potential title area

                        st.write(f"Adding image to slide {i+1} from {image_source['type']} source...") # Debug

                        if image_source["type"] == "url":
                             # _add_image_to_slide handles downloading from URL
                            self._add_image_to_slide(slide, image_source["url"], left=Inches(img_left), top=img_top, width=Inches(img_width_in), height=Inches(img_height_in))

                        elif image_source["type"] == "base64":
                             image_stream = io.BytesIO(base64.b64decode(image_source["data"]))
                             self._add_image_to_slide(slide, image_stream, left=Inches(img_left), top=img_top, width=Inches(img_width_in), height=Inches(img_height_in))


                        # Add caption if provided in main_content
                        caption = "\n".join(content.get("main_content", [])).strip()
                        if caption and not caption.startswith("Error"): # Don't add error messages as captions
                            # Position caption below the image
                            cap_top = img_top + Inches(img_height_in) + Inches(0.1)
                            # Ensure caption box doesn't go below the slide
                            if cap_top + Inches(0.5) < prs.slide_height.inches:
                                cap_box = slide.shapes.add_textbox(Inches(img_left), cap_top, Inches(img_width_in), Inches(0.5))
                                cap_box.text = caption
                                self._format_text_shape(cap_box, font_size=14, alignment=PP_ALIGN.CENTER)
                            else:
                                 st.warning(f"Skipping caption for slide {i+1} as it would exceed slide height.")


                    except Exception as e:
                        st.warning(f"Could not add image to slide {i+1}: {str(e)}")
                        err_box = slide.shapes.add_textbox(Inches(2.5), Inches(3.5), Inches(8), Inches(1))
                        err_box.text = f"(Error adding visual: {str(e)})"
                        self._format_text_shape(err_box, font_size=18, color=RGBColor(255, 0, 0))
                else:
                     st.warning(f"Image slide {i+1} requested but no image source found/generated.")
                     err_box = slide.shapes.add_textbox(Inches(2.5), Inches(3.5), Inches(8), Inches(1))
                     err_box.text = f"(Visual source missing for this slide)"
                     self._format_text_shape(err_box, font_size=18, color=RGBColor(200, 0, 0))

            else: # Handle other slide types (title, content, bullet, conclusion)
                # Set title (using placeholder if available on the chosen layout)
                try:
                    if slide.shapes.title: # Check if title placeholder exists
                        title_shape = slide.shapes.title
                        title_text = content.get("title", slide_title) if slide_type == "title_slide" else slide_title
                        title_shape.text = title_text
                        font_size = 44 if slide_type == "title_slide" else 36
                        self._format_text_shape(title_shape, font_size=font_size, bold=True)
                    else:
                         st.warning(f"Slide {i+1} ('{slide_title}') layout (index {layout_index}) has no title placeholder. Title not set automatically.")
                         # Optionally add a textbox for the title if no placeholder is found
                         # title_shape = slide.shapes.add_textbox(...)

                except AttributeError:
                    # This can happen if the layout doesn't have a shape identified as 'title'
                    st.warning(f"Could not access title shape for slide {i+1} ('{slide_title}') layout (index {layout_index}). Title not set automatically.")

                # Add body content / subtitle
                if slide_type == "title_slide":
                    # Find subtitle placeholder (heuristic approach, often index 1)
                    subtitle_shape = None
                    try:
                        # Check placeholder indices common for subtitles (e.g., 1)
                        if len(slide.placeholders) > 1 and slide.placeholders[1] != slide.shapes.title:
                            subtitle_shape = slide.placeholders[1]
                    except (IndexError, KeyError, AttributeError):
                         pass # Placeholder not found or title check failed

                    if subtitle_shape and subtitle_shape.has_text_frame:
                        subtitle_frame = subtitle_shape.text_frame
                        subtitle_frame.clear() # Clear existing text

                        subtitle_text = content.get("subtitle", "")
                        presenter_text = content.get("presenter", "")
                        date_text = content.get("date", "")

                        p = subtitle_frame.add_paragraph()
                        p.text = subtitle_text
                        self._format_text_shape(p, font_size=28, bold=False)

                        # Add some space
                        p = subtitle_frame.add_paragraph()
                        p.space_after = Pt(12) # Add space after subtitle

                        p = subtitle_frame.add_paragraph()
                        p.text = presenter_text
                        self._format_text_shape(p, font_size=20, bold=False)

                        p = subtitle_frame.add_paragraph()
                        p.text = date_text
                        self._format_text_shape(p, font_size=20, bold=False)

                        self._format_text_shape(subtitle_frame.paragraphs[0], alignment=PP_ALIGN.CENTER) # Center the subtitle text frame content

                    else:
                        st.warning(f"Could not find subtitle placeholder or text frame for title slide {i+1}.")

                elif slide_type in ["content_slide", "bullet_point_slide", "conclusion_slide"]:
                    # Find main content placeholder (heuristic approach)
                    body_shape = self._find_body_placeholder(slide)

                    if body_shape and body_shape.has_text_frame:
                        body_frame = body_shape.text_frame
                        body_frame.clear()
                        body_frame.word_wrap = True

                        main_content_list = content.get("main_content", [])
                        if isinstance(main_content_list, list):
                            main_content_list = [str(point).strip() for point in main_content_list if str(point).strip()] # Clean list
                            for point_idx, point in enumerate(main_content_list):
                                p = body_frame.add_paragraph()
                                p.text = point
                                # Apply bullets only for bullet_point_slide
                                p.level = 0 if slide_type == "bullet_point_slide" else 0
                                font_size = 24 if slide_type == "bullet_point_slide" else 20
                                self._format_text_shape(p, font_size=font_size)
                                if point_idx < len(main_content_list) - 1:
                                     p.space_after = Pt(6)

                        elif isinstance(main_content_list, str) and main_content_list.strip():
                             p = body_frame.add_paragraph()
                             p.text = main_content_list.strip()
                             self._format_text_shape(p, font_size=20)

                        self._auto_fit_text(body_frame)

                    else:
                         st.warning(f"Could not find main content placeholder for slide {i+1} (type: {slide_type}). Content may not appear.")
            # --- END SLIDE CREATION LOGIC ---

        progress_bar.empty()
        st.success("PowerPoint slide creation complete.")

        ppt_bytes = io.BytesIO()
        prs.save(ppt_bytes)
        ppt_bytes.seek(0)

        return ppt_bytes

    # --- Helper methods (_format_text_shape, _add_image_to_slide, _find_body_placeholder, _auto_fit_text) ---
    # (Keep these largely the same, with _add_image_to_slide updated to handle BytesIO)

    def _find_body_placeholder(self, slide):
        """Helper to find the most likely body placeholder on a slide."""
        # (Implementation remains the same)
        body_shape = None
        title_shape_ref = None
        try: title_shape_ref = slide.shapes.title
        except AttributeError: pass

        # Prioritize placeholder index 1 if it exists and is not the title
        try:
            potential_body = slide.placeholders[1]
            is_title = (title_shape_ref == potential_body) if title_shape_ref else False
            if not is_title: body_shape = potential_body
        except (IndexError, KeyError): pass

        # Fallback: find largest non-title placeholder with a text frame
        if body_shape is None:
            largest_area = 0
            for shape in slide.placeholders:
                 if shape.is_placeholder and hasattr(shape, 'text_frame'):
                    if title_shape_ref and shape == title_shape_ref: continue
                    area = shape.width * shape.height
                    if area > largest_area:
                        largest_area = area
                        body_shape = shape
        return body_shape

    def _auto_fit_text(self, text_frame, min_font_size=12):
        """Rudimentary text auto-fitting by reducing font size if text is long."""
        # (Implementation remains the same)
        try:
            text_length = len(text_frame.text)
            length_threshold_medium = 500
            length_threshold_long = 1000

            if text_length > length_threshold_long:
                reduction_factor = 0.8
                target_min_size = min_font_size
            elif text_length > length_threshold_medium:
                reduction_factor = 0.9
                target_min_size = min_font_size + 2
            else:
                return

            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if hasattr(run.font, 'size') and run.font.size:
                        current_size_pt = run.font.size.pt
                        new_size = max(target_min_size, int(current_size_pt * reduction_factor))
                        run.font.size = Pt(new_size)
        except Exception as e:
             st.warning(f"Minor error during text auto-fit attempt: {e}")


    def _format_text_shape(self, shape, font_size=None, bold=False, color=None, alignment=None):
        """Apply text formatting to paragraphs/runs within a shape or directly to a paragraph/run."""
        # (Implementation remains the same)
        try:
            target = None
            if hasattr(shape, 'text_frame'): target = shape.text_frame
            elif hasattr(shape, 'font'): target = shape # Paragraph or Run

            if target:
                if hasattr(target, 'paragraphs'): # Text Frame
                    for paragraph in target.paragraphs:
                        if alignment is not None: paragraph.alignment = alignment
                        for run in paragraph.runs:
                            if font_size is not None: run.font.size = Pt(font_size)
                            run.font.bold = bold
                            if color is not None: run.font.color.rgb = color
                elif hasattr(target, 'runs'): # Paragraph
                     paragraph = target
                     if alignment is not None: paragraph.alignment = alignment
                     for run in paragraph.runs:
                        if font_size is not None: run.font.size = Pt(font_size)
                        run.font.bold = bold
                        if color is not None: run.font.color.rgb = color
                elif hasattr(target, 'font'): # Run
                    run = target
                    if font_size is not None: run.font.size = Pt(font_size)
                    run.font.bold = bold
                    if color is not None: run.font.color.rgb = color
        except AttributeError: pass
        except Exception as e: st.warning(f"Minor formatting error ignored: {e}")


    def _add_image_to_slide(self, slide, image_source, left=None, top=None, width=None, height=None):
        """
        Add an image from a URL or BytesIO stream to a slide.
        image_source can be a URL string or a BytesIO object.
        """
        st.write(f"Attempting to add image from source type: {type(image_source)}")
        image_stream = None
        source_desc = "Unknown"

        try:
            if isinstance(image_source, str): # It's a URL
                source_desc = image_source
                headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/90.0.4430.93 Safari/537.36'}
                response = requests.get(image_source, headers=headers, timeout=60, stream=True, allow_redirects=True)
                response.raise_for_status()

                content_type = response.headers.get('Content-Type', '').lower()
                if 'image' not in content_type:
                    st.warning(f"URL {image_source} returned content type '{content_type}', which might not be an image.")

                image_stream = io.BytesIO()
                size = 0
                max_size = 20 * 1024 * 1024 # 20 MB limit
                for chunk in response.iter_content(chunk_size=8192):
                    size += len(chunk)
                    if size > max_size:
                        raise ValueError(f"Image download exceeded max size limit ({max_size/1024/1024} MB)")
                    image_stream.write(chunk)
                image_stream.seek(0)

                if image_stream.getbuffer().nbytes == 0:
                    raise ValueError("Image download resulted in empty content.")

                st.write("Image downloaded successfully.")

            elif isinstance(image_source, io.BytesIO): # It's already a BytesIO stream (e.g., from local file)
                 image_stream = image_source
                 image_stream.seek(0) # Ensure stream is at the beginning
                 source_desc = "BytesIO stream (local file or base64)"
                 st.write("Using provided BytesIO image stream.")

            else:
                 raise TypeError("image_source must be a URL string or BytesIO object")


            if left is None: left = Inches(1.0)
            if top is None: top = Inches(1.0)

            if width is not None and height is not None:
                pic = slide.shapes.add_picture(image_stream, left, top, width=width, height=height)
            elif width is not None:
                pic = slide.shapes.add_picture(image_stream, left, top, width=width)
            elif height is not None:
                 pic = slide.shapes.add_picture(image_stream, left, top, height=height)
            else:
                 pic = slide.shapes.add_picture(image_stream, left, top)
            st.write(f"Image added to slide (Shape ID: {pic.shape_id}).")

        except requests.exceptions.Timeout:
             st.error(f"Timeout error downloading image: {source_desc}")
             raise ValueError(f"Timeout downloading image: {source_desc}")
        except requests.exceptions.RequestException as req_err:
             st.error(f"Network error downloading image: {req_err} (Source: {source_desc})")
             if hasattr(req_err, 'response') and req_err.response is not None:
                 st.error(f"Status Code: {req_err.response.status_code}")
             raise ValueError(f"Failed to download image: {req_err}")
        except FileNotFoundError as fnf_err:
             st.error(f"File not found error adding image: {fnf_err} (Source: {source_desc})")
             raise ValueError(f"File not found: {fnf_err}")
        except Exception as e:
            st.error(f"Failed to add image to slide from source ({source_desc}): {str(e)}")
            raise ValueError(f"Failed to process or add image from source ({source_desc}): {str(e)}")